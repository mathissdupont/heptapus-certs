"""PowerPoint rendering for HeptaCert presentation decks."""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _hex_to_rgb(value: str, fallback: str = "#2563eb") -> RGBColor:
    raw = (value or fallback).strip().lstrip("#")
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    try:
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except ValueError:
        raw = fallback.lstrip("#")
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _text(frame: Any, value: str, *, size: int, bold: bool = False, color: str = "#0f172a") -> None:
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(color, color)


def _add_bullets(frame: Any, bullets: list[str], *, color: str) -> None:
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0
        paragraph.font.size = Pt(21)
        paragraph.font.color.rgb = _hex_to_rgb(color, color)
        paragraph.space_after = Pt(8)


def render_deck_pptx(deck: Any) -> bytes:
    """Render a deck row or deck-like object to PPTX bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = deck.theme or {}
    slides = deck.slides or []
    primary = str(theme.get("primary") or "#2563eb")
    background = str(theme.get("background") or "#f8fafc")
    foreground = str(theme.get("foreground") or "#0f172a")

    if not slides:
        slides = [{"title": deck.title, "subtitle": deck.description or "", "layout": "title"}]

    for index, slide_data in enumerate(slides):
        layout = str(slide_data.get("layout") or ("title" if index == 0 else "bullets"))
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(str(slide_data.get("background") or background), background)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _hex_to_rgb(primary, primary)
        accent.line.fill.background()

        if layout == "title":
            title_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.15), Inches(11.7), Inches(1.0))
            _text(title_box.text_frame, str(slide_data.get("title") or deck.title), size=40, bold=True, color=foreground)
            subtitle = str(slide_data.get("subtitle") or slide_data.get("notes") or deck.description or "")
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.88), Inches(3.35), Inches(10.9), Inches(0.75))
                _text(sub_box.text_frame, subtitle, size=20, color="#475569")
        else:
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.9), Inches(0.75))
            _text(title_box.text_frame, str(slide_data.get("title") or f"Slide {index + 1}"), size=30, bold=True, color=foreground)
            bullets = slide_data.get("bullets") if isinstance(slide_data.get("bullets"), list) else []
            if not bullets and slide_data.get("body"):
                bullets = [str(slide_data.get("body"))]
            body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.2), Inches(4.8))
            _add_bullets(body_box.text_frame, [str(item) for item in bullets[:8]], color=foreground)

        footer = slide.shapes.add_textbox(Inches(0.78), Inches(6.95), Inches(11.8), Inches(0.25))
        _text(footer.text_frame, f"HeptaCert · {index + 1}", size=9, color="#64748b")

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()
